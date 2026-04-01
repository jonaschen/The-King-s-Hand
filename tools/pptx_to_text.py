#!/usr/bin/env python3
"""
pptx_to_text.py — Convert a PowerPoint (.pptx) file to structured plain text
formatted for submission to The King's Hand analysis system.

Output format matches TC-015 (project_eon_ai_slides.md): each slide is prefixed
with [SLIDE N], tables are formatted as aligned text, speaker notes are optional.

USAGE:
    python3 pptx_to_text.py presentation.pptx
    python3 pptx_to_text.py presentation.pptx --notes
    python3 pptx_to_text.py presentation.pptx --output extracted.txt
    python3 pptx_to_text.py presentation.pptx > extracted.txt

REQUIRES:
    pip install python-pptx
"""

import sys
import argparse
from pathlib import Path


def format_table(table_shape):
    """Format a pptx table as aligned plain text."""
    rows = []
    for row in table_shape.table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    # Deduplicate merged cells (pptx repeats merged cell text across spans)
    col_count = max(len(r) for r in rows)
    col_widths = []
    for c in range(col_count):
        width = max((len(r[c]) if c < len(r) else 0) for r in rows)
        col_widths.append(max(width, 3))

    lines = []
    for i, row in enumerate(rows):
        padded = []
        for c in range(col_count):
            cell = row[c] if c < len(row) else ""
            padded.append(cell.ljust(col_widths[c]))
        lines.append("  " + "   ".join(padded).rstrip())
        # Separator after header row
        if i == 0:
            lines.append("  " + "   ".join("─" * col_widths[c] for c in range(col_count)))

    return "\n".join(lines)


def pptx_to_text(path, include_notes=False):
    """
    Extract text from a .pptx file, slide by slide.

    Returns a string in [SLIDE N] format suitable for pasting into
    The King's Hand agent session.
    """
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(path)
    output_sections = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = [f"[SLIDE {slide_num}]"]

        # Separate title placeholders from body content
        title_text = None
        body_parts = []

        for shape in slide.shapes:
            # Skip shapes without text or tables
            if not shape.has_text_frame and not shape.has_table:
                continue

            if shape.has_table:
                table_text = format_table(shape)
                if table_text:
                    body_parts.append(table_text)
                continue

            # Text frame
            text = shape.text_frame.text.strip()
            if not text:
                continue

            # Classify by placeholder type
            is_title = False
            try:
                ph_fmt = shape.placeholder_format
                if ph_fmt is not None and ph_fmt.idx == 0:
                    is_title = True
            except ValueError:
                pass  # shape.placeholder_format raises ValueError for non-placeholders

            if is_title:
                title_text = text
            else:
                body_parts.append(text)

        # Emit title first
        if title_text:
            slide_lines.append(title_text)

        slide_lines.extend(body_parts)

        # Speaker notes (optional)
        if include_notes:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip() if notes_frame else ""
                if notes_text:
                    slide_lines.append(f"\n[SPEAKER NOTES: {notes_text}]")
            except Exception:
                pass  # notes not available on this slide

        output_sections.append("\n".join(slide_lines))

    return "\n\n".join(output_sections)


def main():
    parser = argparse.ArgumentParser(
        description="Convert .pptx to structured text for The King's Hand"
    )
    parser.add_argument("file", help="Path to the .pptx file")
    parser.add_argument(
        "--notes", action="store_true",
        help="Include speaker notes (if present)"
    )
    parser.add_argument(
        "--output", "-o", metavar="FILE",
        help="Write output to FILE instead of stdout"
    )
    args = parser.parse_args()

    # Dependency check
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print(
            "ERROR: python-pptx is not installed.\n"
            "Run:   pip install python-pptx",
            file=sys.stderr
        )
        sys.exit(1)

    path = Path(args.file)
    if not path.exists():
        print(f"ERROR: File not found: {path}", file=sys.stderr)
        sys.exit(1)
    if path.suffix.lower() != ".pptx":
        print(
            f"WARNING: File extension is '{path.suffix}', expected .pptx. Proceeding anyway.",
            file=sys.stderr
        )

    text = pptx_to_text(path, include_notes=args.notes)

    if args.output:
        out_path = Path(args.output)
        out_path.write_text(text, encoding="utf-8")
        print(f"Written to: {out_path}", file=sys.stderr)
    else:
        print(text)


if __name__ == "__main__":
    main()
